"""
elec_week_renderer.py
电量周度数据 PNG 渲染器
方案：填充 PPT 模板 → LibreOffice 转 PNG，保证与模板完全一致
"""

import os
import re
import shutil
import subprocess
import tempfile
from datetime import datetime
from copy import deepcopy

from pptx import Presentation
from pptx.util import Pt
from lxml import etree

_HERE        = os.path.dirname(os.path.abspath(__file__))
TEMPLATE_PATH = os.path.join(_HERE, "assets", "电量周度数据_模板.pptx")
OUT_DIR       = os.path.join(_HERE, "output")


def _set_cell_text(cell, text: str):
    """替换单元格文字，保留原有第一个 run 的字体格式"""
    tf = cell.text_frame
    p_ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    txBody = tf._txBody

    # 取第一个段落第一个 run 的格式
    first_para = tf.paragraphs[0]
    ref_rPr = None
    if first_para.runs:
        ref_rPr = first_para.runs[0]._r.find(f'{{{p_ns}}}rPr')

    # 清空所有段落内容（保留段落元素本身）
    for p_el in txBody.findall(f'{{{p_ns}}}p'):
        for r_el in p_el.findall(f'{{{p_ns}}}r'):
            p_el.remove(r_el)

    # 只写第一个段落
    p_el = txBody.findall(f'{{{p_ns}}}p')[0]
    r_el = etree.SubElement(p_el, f'{{{p_ns}}}r')
    if ref_rPr is not None:
        r_el.insert(0, deepcopy(ref_rPr))
    t_el = etree.SubElement(r_el, f'{{{p_ns}}}t')
    t_el.text = text


def _set_cell_two_lines(cell, line1: str, line2: str):
    """
    设置单元格两行文字：
    - 第一行：line1（居中）
    - 第二行：line2（居中，自动缩小字号确保单行）
    两行都保留原始 rPr 格式
    """
    tf = cell.text_frame
    p_ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    txBody = tf._txBody

    # 取第一个段落第一个 run 的格式作为模板
    first_para = tf.paragraphs[0]
    ref_rPr = None
    if first_para.runs:
        ref_rPr = first_para.runs[0]._r.find(f'{{{p_ns}}}rPr')

    # 取第一个段落的 pPr（对齐等段落格式）作为模板
    ref_pPr = first_para._p.find(f'{{{p_ns}}}pPr')

    # 清空所有段落
    for p_el in txBody.findall(f'{{{p_ns}}}p'):
        txBody.remove(p_el)

    def _make_para(text, rPr_template, pPr_template, font_size_pt=None):
        p_el = etree.SubElement(txBody, f'{{{p_ns}}}p')
        if pPr_template is not None:
            p_el.insert(0, deepcopy(pPr_template))
        r_el = etree.SubElement(p_el, f'{{{p_ns}}}r')
        rPr = deepcopy(rPr_template) if rPr_template is not None else etree.Element(f'{{{p_ns}}}rPr')
        # 如果需要缩小字号
        if font_size_pt is not None:
            rPr.set('sz', str(int(font_size_pt * 100)))
        r_el.insert(0, rPr)
        t_el = etree.SubElement(r_el, f'{{{p_ns}}}t')
        t_el.text = text
        return p_el

    _make_para(line1, ref_rPr, ref_pPr)
    # 第二行日期：字号缩小到 32pt（原始 40pt），确保长日期单行显示
    _make_para(line2, ref_rPr, ref_pPr, font_size_pt=32)


def _set_textbox_text(shape, text: str):
    """替换文本框整段文字，保留格式"""
    tf = shape.text_frame
    p_ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    txBody = tf._txBody
    paras = txBody.findall(f'{{{p_ns}}}p')
    if not paras:
        return
    first_p = paras[0]
    runs = first_p.findall(f'{{{p_ns}}}r')
    ref_rPr = None
    if runs:
        ref_rPr = runs[0].find(f'{{{p_ns}}}rPr')
        for r in runs:
            first_p.remove(r)

    r_el = etree.SubElement(first_p, f'{{{p_ns}}}r')
    if ref_rPr is not None:
        r_el.insert(0, deepcopy(ref_rPr))
    t_el = etree.SubElement(r_el, f'{{{p_ns}}}t')
    t_el.text = text


def fill_template(data: dict, out_pptx: str):
    prs  = Presentation(TEMPLATE_PATH)
    slide = prs.slides[0]

    week_range  = data.get("week_range",  "")
    report_date = data.get("report_date", "")
    region      = data.get("region",      "公司经营区")
    month_range = data.get("month_range", "")
    year_range  = data.get("year_range",  "")
    notes_text  = data.get("notes",       "")
    week  = data.get("week",  {})
    month = data.get("month", {})
    year  = data.get("year",  {})

    for shape in slide.shapes:
        name = shape.name

        # 日期范围 TextBox3（top≈95mm）
        if name == "TextBox 3" and shape.top / 914400 * 25.4 > 90:
            _set_textbox_text(shape, f"({week_range})")

        # 地区/报送 TextBox4（top≈130mm）
        if name == "TextBox 4" and shape.top / 914400 * 25.4 > 125 and shape.top / 914400 * 25.4 < 145:
            p_ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
            tf = shape.text_frame
            txBody = tf._txBody
            paras = txBody.findall(f'{{{p_ns}}}p')
            ref_rPr = None
            if paras and tf.paragraphs[0].runs:
                ref_rPr = tf.paragraphs[0].runs[0]._r.find(f'{{{p_ns}}}rPr')
            for p_el in paras:
                txBody.remove(p_el)

            left_txt  = f"地区范围：{region}"
            right_txt = f"报送日期：{report_date}"

            # 一个段落，左侧 run + 右侧 run，段落对齐用 'thaiDist' 无法实现
            # 最终方案：写两行，第一行左对齐，第二行右对齐，然后把文本框高度缩到只有一行
            # 但更简单：段落左对齐写 left_txt，用 tab + 右缩进写 right_txt
            # 实际最稳定：一段，left + 大量空格 + right，空格用 &nbsp; 等宽字符
            PAD_EMU = int(8 * 36000)  # 8mm 两侧留白
            p_el = etree.SubElement(txBody, f'{{{p_ns}}}p')
            pPr_el = etree.SubElement(p_el, f'{{{p_ns}}}pPr')
            pPr_el.set('algn', 'l')

            tab_w_emu = shape.width
            tabLst = etree.SubElement(pPr_el, f'{{{p_ns}}}tabLst')
            tab_el = etree.SubElement(tabLst, f'{{{p_ns}}}tab')
            tab_el.set('pos', str(int(tab_w_emu - PAD_EMU)))  # 右侧留白
            tab_el.set('algn', 'r')

            # 左侧文字（前加空格留白）
            r_left = etree.SubElement(p_el, f'{{{p_ns}}}r')
            if ref_rPr is not None:
                r_left.insert(0, deepcopy(ref_rPr))
            t_left = etree.SubElement(r_left, f'{{{p_ns}}}t')
            t_left.text = '  ' + left_txt  # 两个空格作为左侧留白

            # 右侧文字：用 fldType=right tab stop 实现右对齐
            # 最简单方式：在 pPr 里加右 tabStop，然后插入 tab
            # tabStop at 文本框宽度
            tab_w_emu = shape.width
            tabLst = etree.SubElement(pPr_el, f'{{{p_ns}}}tabLst')
            tab_el = etree.SubElement(tabLst, f'{{{p_ns}}}tab')
            tab_el.set('pos', str(int(tab_w_emu - PAD_EMU)))  # 右侧留白
            tab_el.set('algn', 'r')

            # tab 字符
            r_tab = etree.SubElement(p_el, f'{{{p_ns}}}r')
            if ref_rPr is not None:
                r_tab.insert(0, deepcopy(ref_rPr))
            t_tab = etree.SubElement(r_tab, f'{{{p_ns}}}t')
            t_tab.text = '\t'

            # 右侧文字
            r_right = etree.SubElement(p_el, f'{{{p_ns}}}r')
            if ref_rPr is not None:
                r_right.insert(0, deepcopy(ref_rPr))
            t_right = etree.SubElement(r_right, f'{{{p_ns}}}t')
            t_right.text = right_txt

        # 备注 TextBox3（top≈773mm）
        if name == "TextBox 3" and shape.top / 914400 * 25.4 > 760:
            _set_textbox_text(shape, f"备注：{notes_text}")

        # 表格
        if hasattr(shape, "table"):
            tbl = shape.table

            # 周期列：名称+日期各一行，日期字号缩小确保单行
            if month_range:
                _set_cell_two_lines(tbl.cell(3, 0), "月累计电量", f"({month_range})")
            else:
                _set_cell_text(tbl.cell(3, 0), "月累计电量")

            if year_range:
                _set_cell_two_lines(tbl.cell(5, 0), "年累计电量", f"({year_range})")
            else:
                _set_cell_text(tbl.cell(5, 0), "年累计电量")

            # 数值/同比
            _set_cell_text(tbl.cell(1, 2), week.get("collect", ""))
            _set_cell_text(tbl.cell(1, 3), week.get("collect_yoy", ""))
            _set_cell_text(tbl.cell(2, 2), week.get("sale", ""))
            _set_cell_text(tbl.cell(2, 3), week.get("sale_yoy", ""))

            _set_cell_text(tbl.cell(3, 2), month.get("collect", ""))
            _set_cell_text(tbl.cell(3, 3), month.get("collect_yoy", ""))
            _set_cell_text(tbl.cell(4, 2), month.get("sale", ""))
            _set_cell_text(tbl.cell(4, 3), month.get("sale_yoy", ""))

            _set_cell_text(tbl.cell(5, 2), year.get("collect", ""))
            _set_cell_text(tbl.cell(5, 3), year.get("collect_yoy", ""))
            _set_cell_text(tbl.cell(6, 2), year.get("sale", ""))
            _set_cell_text(tbl.cell(6, 3), year.get("sale_yoy", ""))

    prs.save(out_pptx)


def render_elec_week_png(data: dict, output_path: str = None) -> str:
    # 确保 LibreOffice 能找到微软雅黑字体
    fonts_dir = os.path.join(_HERE, "fonts")
    lo_fonts  = os.path.expanduser("~/.fonts")
    if os.path.isdir(fonts_dir) and not os.path.exists(os.path.join(lo_fonts, "msyh.ttf")):
        os.makedirs(lo_fonts, exist_ok=True)
        import shutil as _sh
        for f in os.listdir(fonts_dir):
            if f.endswith((".ttf", ".otf")):
                _sh.copy(os.path.join(fonts_dir, f), lo_fonts)
        os.system("fc-cache -f ~/.fonts 2>/dev/null")
    # 固定字段默认值（不从文字中解析，直接用模板里的值）
    defaults = {
        "week_range":  "",
        "report_date": "",
        "region":      "公司经营区",
        "month_range": "",
        "year_range":  "",
    }
    # parser 只返回 week/month/year/notes，合并固定字段
    merged = {**defaults, **data}

    os.makedirs(OUT_DIR, exist_ok=True)
    if not output_path:
        ts = datetime.now().strftime('%Y-%m-%d')
        output_path = os.path.join(OUT_DIR, f"elec_week_{ts}.png")

    with tempfile.TemporaryDirectory() as tmpdir:
        tmp_pptx = os.path.join(tmpdir, "filled.pptx")
        fill_template(merged, tmp_pptx)

        # 找 LibreOffice 可执行文件（不同系统命令名不同）
        import shutil as _shutil
        lo_cmd = (_shutil.which("libreoffice") or
                  _shutil.which("soffice") or
                  "/usr/bin/libreoffice")

        # LibreOffice 转 PNG
        result = subprocess.run(
            [lo_cmd, "--headless", "--convert-to", "png",
             "--outdir", tmpdir, tmp_pptx],
            capture_output=True, text=True, timeout=120
        )
        print(result.stdout)
        if result.returncode != 0:
            print(result.stderr)
            raise RuntimeError(f"LibreOffice 转换失败: {result.stderr}")

        # LibreOffice 输出文件名是 filled.png
        tmp_png = os.path.join(tmpdir, "filled.png")
        if not os.path.exists(tmp_png):
            # 找一下实际输出
            pngs = [f for f in os.listdir(tmpdir) if f.endswith(".png")]
            if not pngs:
                raise RuntimeError(f"未找到输出 PNG，tmpdir内容: {os.listdir(tmpdir)}")
            tmp_png = os.path.join(tmpdir, pngs[0])

        shutil.copy(tmp_png, output_path)

    print(f"✅ 已生成：{output_path}")
    return output_path


if __name__ == "__main__":
    MOCK = {
        "week_range":  "2026年6月29日-2026年7月5日",
        "report_date": "2026年7月6日",
        "region":      "公司经营区",
        "month_range": "7月1日-5日",
        "year_range":  "1月1日-7月5日",
        "week":  {"collect":"1490.34","collect_yoy":"-2.8%","sale":"1376.47","sale_yoy":"-2.6%"},
        "month": {"collect":"1073.01","collect_yoy":"-5.0%","sale":"990.93", "sale_yoy":"-4.9%"},
        "year":  {"collect":"35296.88","collect_yoy":"4.4%","sale":"32532.54","sale_yoy":"4.4%"},
        "notes": "一是本周电量负增长，主要受江浙等地降雨天气影响，平均最高气温较同期低4.8℃，居民用电量同比下降15.4%。二是周采集电量增速低于售电量增速，主要由于江苏、浙江、湖南、山东、上海、安徽等省份连续阴雨天气，分布式光伏出力下降，自用电量减少。",
    }
    render_elec_week_png(MOCK, output_path="output/elec_week_test.png")
